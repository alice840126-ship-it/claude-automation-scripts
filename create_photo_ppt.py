#!/usr/bin/env python3
"""
사진을 A4 PPT로 자동 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys
from PIL import Image
import io
import tempfile
import os

# A4 크기 (인치)
A4_WIDTH = 8.27
A4_HEIGHT = 11.69

def convert_to_jpeg(img_path):
    """이미지를 JPEG로 변환"""
    try:
        img = Image.open(img_path)

        # RGB로 변환 (필요시)
        if img.mode != 'RGB':
            img = img.convert('RGB')

        # 임시 파일로 저장
        temp_fd, temp_path = tempfile.mkstemp(suffix='.jpg')
        os.close(temp_fd)

        img.save(temp_path, 'JPEG', quality=95)
        return temp_path
    except Exception as e:
        print(f"⚠️ 이미지 변환 실패: {img_path.name} - {e}")
        return None

def create_photo_ppt(image_folder, output_path):
    """사진을 A4 PPT로 생성"""

    # 이미지 파일들 가져오기
    folder = Path(image_folder)
    images = sorted(folder.glob("*.JPG")) + sorted(folder.glob("*.jpg")) + sorted(folder.glob("*.PNG")) + sorted(folder.glob("*.png"))

    if not images:
        print("❌ 이미지 파일 없음")
        return

    print(f"📷 {len(images)}장의 사진 발견")

    # 프레젠테이션 생성
    prs = Presentation()

    # 페이지 크기를 A4로 설정
    prs.slide_width = Inches(A4_WIDTH)
    prs.slide_height = Inches(A4_HEIGHT)

    # 빈 슬라이드 추가
    blank_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_slide_layout)

    # 그리드 설정 (4행 x 2열)
    rows = 4
    cols = 2

    # 여백
    margin_x = 0.3  # 왼쪽/오른쪽 여백
    margin_y = 0.5  # 위/아래 여백

    # 그림 영역 계산
    pic_width = (A4_WIDTH - margin_x * 3) / 2
    pic_height = (A4_HEIGHT - margin_y * (rows + 1)) / rows

    # 사진 배치
    temp_files = []  # 임시 파일 추적

    for idx, img_path in enumerate(images[:8]):  # 최대 8장
        row = idx // cols
        col = idx % cols

        # 위치 계산
        left = margin_x + col * (pic_width + margin_x)
        top = margin_y + row * (pic_height + margin_y)

        # 이미지를 JPEG로 변환
        jpeg_path = convert_to_jpeg(img_path)
        if not jpeg_path:
            continue

        temp_files.append(jpeg_path)

        # 이미지 추가 (비율 유지하며 적합)
        pic = slide.shapes.add_picture(jpeg_path, left, top, width=Inches(pic_width))

        # 이미지가 높이를 초과하면 조정
        max_height = Inches(pic_height)
        if pic.height > max_height:
            # 높이에 맞춰서 비율 조정
            ratio = max_height.inches / pic.height.inches
            new_width = pic.width.inches * ratio
            pic.height = max_height
            pic.width = Inches(new_width)
            # 가운데 정렬
            pic.left = Inches(left + (pic_width - new_width) / 2)

        # 파일명 표시
        textbox = slide.shapes.add_textbox(
            left,
            top + pic_height + 0.05,
            Inches(pic_width),
            Inches(0.2)
        )
        text_frame = textbox.text_frame
        text_frame.text = img_path.name
        text_frame.word_wrap = True

        # 폰트 설정
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.name = "맑은 고딕"

    # 저장
    prs.save(output_path)

    # 임시 파일 삭제
    for temp_file in temp_files:
        try:
            os.unlink(temp_file)
        except:
            pass

    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"   크기: {A4_WIDTH:.2f}\" x {A4_HEIGHT:.2f}\" (A4)")
    print(f"   사진: {len(images[:8])}장")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python3 create_photo_ppt.py <이미지_폴더>")
        sys.exit(1)

    image_folder = sys.argv[1]
    output_path = str(Path(image_folder).parent / "사진_배열.pptx")

    create_photo_ppt(image_folder, output_path)
